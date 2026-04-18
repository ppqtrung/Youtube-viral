"""Generate YouTube Viral Advisor PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette (matches app) ─────────────────────────────────────────────────────
SLATE_900 = RGBColor(0x0f, 0x17, 0x2a)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_600 = RGBColor(0x47, 0x55, 0x69)
SLATE_500 = RGBColor(0x64, 0x74, 0x8b)
SLATE_400 = RGBColor(0x94, 0xa3, 0xb8)
SLATE_200 = RGBColor(0xe2, 0xe8, 0xf0)
SLATE_100 = RGBColor(0xf1, 0xf5, 0xf9)
SLATE_50  = RGBColor(0xf8, 0xfa, 0xfc)
WHITE     = RGBColor(0xff, 0xff, 0xff)
ACCENT    = RGBColor(0x25, 0x63, 0xeb)   # blue-600
ACCENT_D  = RGBColor(0x1d, 0x4e, 0xd8)  # blue-700
SUCCESS   = RGBColor(0x05, 0x96, 0x69)
DANGER    = RGBColor(0xdc, 0x26, 0x26)
WARNING   = RGBColor(0xd9, 0x77, 0x06)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # completely blank


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor | None = None,
             line: RGBColor | None = None, line_w: int = 0):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_w:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text: str, x, y, w, h,
             size: int = 20, bold: bool = False, color: RGBColor = SLATE_900,
             align=PP_ALIGN.LEFT, wrap: bool = True, italic: bool = False) -> None:
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_para_box(slide, lines: list[tuple[str, int, bool, RGBColor]],
                 x, y, w, h, line_spacing: float = 1.3):
    """Multi-paragraph text box. Each line = (text, size, bold, color)."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"


def bullet_box(slide, items: list[str], x, y, w, h,
               size: int = 18, color: RGBColor = SLATE_700,
               accent: RGBColor = ACCENT, dot: str = "•"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Bullet dot
        r1 = p.add_run()
        r1.text = dot + "  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = accent
        r1.font.name = "Calibri"
        # Text
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = "Calibri"


def kpi_card(slide, label: str, value: str, sub: str,
             x, y, w=Inches(2.8), h=Inches(1.45),
             val_color: RGBColor = SLATE_900):
    add_rect(slide, x, y, w, h, fill=WHITE, line=SLATE_200, line_w=8000)
    add_text(slide, label.upper(), x + Inches(0.18), y + Inches(0.15),
             w - Inches(0.3), Inches(0.3), size=8, color=SLATE_500, bold=True)
    add_text(slide, value, x + Inches(0.15), y + Inches(0.42),
             w - Inches(0.3), Inches(0.65), size=28, bold=True, color=val_color)
    if sub:
        add_text(slide, sub, x + Inches(0.18), y + Inches(1.1),
                 w - Inches(0.3), Inches(0.3), size=9, color=SLATE_500)


def section_header(slide, number: str, title: str,
                   x=Inches(0.6), y=Inches(0.45)):
    add_text(slide, number, x, y, Inches(1), Inches(0.3),
             size=9, color=SLATE_400, bold=True)
    add_rect(slide, x, y + Inches(0.35), Inches(0.04), Inches(0.38), fill=ACCENT)
    add_text(slide, title, x + Inches(0.18), y + Inches(0.3),
             Inches(10), Inches(0.5), size=28, bold=True, color=SLATE_900)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title + team."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)

    # Left accent bar
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)

    # Top label
    add_text(s, "HACKATHON 2026  |  DATA SCIENCE TRACK",
             Inches(0.25), Inches(0.38), Inches(9), Inches(0.35),
             size=9, color=SLATE_400, bold=True)

    # Main title
    add_text(s, "YouTube Viral Advisor",
             Inches(0.25), Inches(1.6), Inches(9), Inches(1.4),
             size=52, bold=True, color=SLATE_900)

    # Subtitle
    add_text(s,
             "Predicting video virality before upload using\n"
             "gradient-boosted ML, NLP, and 178,000 trending videos.",
             Inches(0.25), Inches(3.15), Inches(8.5), Inches(1.1),
             size=20, color=SLATE_600)

    # Divider
    add_rect(s, Inches(0.25), Inches(4.45), Inches(2.5), Inches(0.025), fill=ACCENT)

    # Team box on right
    add_rect(s, Inches(9.5), Inches(1.4), Inches(3.5), Inches(4.8),
             fill=SLATE_50, line=SLATE_200, line_w=6000)
    add_text(s, "TEAM", Inches(9.7), Inches(1.6), Inches(3), Inches(0.35),
             size=9, color=SLATE_400, bold=True)
    add_text(s, "Team Name", Inches(9.7), Inches(2.0), Inches(3), Inches(0.45),
             size=18, bold=True, color=SLATE_900)
    add_text(s,
             "Member 1\nMember 2\nMember 3",
             Inches(9.7), Inches(2.55), Inches(3), Inches(1.2),
             size=15, color=SLATE_700)
    add_text(s, "Submission includes:", Inches(9.7), Inches(3.9),
             Inches(3), Inches(0.35), size=10, color=SLATE_500, bold=True)
    add_text(s,
             "• Video presentation\n• This slide deck\n• Full source code (app.py,\n  analysis.py, model.py)",
             Inches(9.7), Inches(4.3), Inches(3.1), Inches(1.5),
             size=12, color=SLATE_600)

    # Dataset badge
    add_rect(s, Inches(0.25), Inches(4.9), Inches(5.2), Inches(0.55),
             fill=SLATE_50, line=SLATE_200, line_w=6000)
    add_text(s, "178,267 videos  |  11 countries  |  15 categories  |  2026 dataset",
             Inches(0.45), Inches(5.0), Inches(4.8), Inches(0.38),
             size=11, color=SLATE_600)


def slide_problem(prs):
    """Slide 2 — Problem statement."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "01", "Problem Statement")

    add_text(s,
             "Small YouTubers upload blind.",
             Inches(0.6), Inches(1.45), Inches(9), Inches(0.7),
             size=32, bold=True, color=SLATE_900)

    add_text(s,
             "Before a video goes live, creators have no data-driven signal on whether\n"
             "their choices — title, timing, tags, category — match what actually trends.",
             Inches(0.6), Inches(2.3), Inches(10.5), Inches(0.85),
             size=18, color=SLATE_600)

    # Three pain cards
    cards = [
        ("Wrong upload time",   "Peak hours matter. Most creators publish whenever is convenient."),
        ("Poor title strategy", "Title length and keywords directly impact click-through rate."),
        ("Tag guesswork",       "Too few or too many tags reduce discoverability in recommendations."),
    ]
    cx = Inches(0.6)
    for label, desc in cards:
        add_rect(s, cx, Inches(3.45), Inches(3.85), Inches(2.6),
                 fill=SLATE_50, line=SLATE_200, line_w=8000)
        add_rect(s, cx, Inches(3.45), Inches(0.055), Inches(2.6), fill=DANGER)
        add_text(s, label, cx + Inches(0.2), Inches(3.65), Inches(3.4), Inches(0.5),
                 size=16, bold=True, color=SLATE_900)
        add_text(s, desc, cx + Inches(0.2), Inches(4.25), Inches(3.45), Inches(1.4),
                 size=13, color=SLATE_600)
        cx += Inches(4.1)


def slide_solution(prs):
    """Slide 3 — Solution overview."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "02", "Solution Overview")

    # Flow boxes
    steps = [
        ("Raw Data",     "11 country CSVs\n178K+ videos",          SLATE_100, SLATE_900),
        ("Feature Eng.", "NLP + timing\n+ engagement ratios",       SLATE_100, SLATE_900),
        ("ML Model",     "HistGradientBoosting\n+ Calibration",     ACCENT,    WHITE),
        ("Virality Score","0-100% calibrated\nprobability",         SLATE_900, WHITE),
    ]
    arrow_x = Inches(0.45)
    for i, (title, body, bg, fg) in enumerate(steps):
        bx = arrow_x + Inches(0.1)
        add_rect(s, bx, Inches(1.8), Inches(2.7), Inches(1.8),
                 fill=bg, line=SLATE_200, line_w=8000)
        add_text(s, title, bx + Inches(0.18), Inches(1.98),
                 Inches(2.35), Inches(0.4), size=14, bold=True, color=fg)
        add_text(s, body, bx + Inches(0.18), Inches(2.45),
                 Inches(2.35), Inches(0.9), size=12, color=fg if bg == ACCENT or bg == SLATE_900 else SLATE_600)
        if i < len(steps) - 1:
            ax = bx + Inches(2.75)
            add_text(s, "→", ax, Inches(2.35), Inches(0.5), Inches(0.5),
                     size=22, bold=True, color=SLATE_400, align=PP_ALIGN.CENTER)
        arrow_x += Inches(3.05)

    # Three output features
    add_text(s, "Three outputs per prediction", Inches(0.6), Inches(4.0),
             Inches(9), Inches(0.45), size=17, bold=True, color=SLATE_900)
    outputs = [
        ("Virality Score",     "Calibrated probability (0-100%) of reaching top-25% views"),
        ("Benchmark Table",    "Side-by-side comparison of inputs vs. viral video benchmarks"),
        ("Recommendations",    "Data-driven suggestions + NLP viral keyword hints"),
    ]
    bx = Inches(0.6)
    for title, desc in outputs:
        add_rect(s, bx, Inches(4.6), Inches(3.95), Inches(2.35),
                 fill=WHITE, line=SLATE_200, line_w=8000)
        add_rect(s, bx, Inches(4.6), Inches(3.95), Inches(0.06), fill=ACCENT)
        add_text(s, title, bx + Inches(0.18), Inches(4.78), Inches(3.6), Inches(0.45),
                 size=14, bold=True, color=SLATE_900)
        add_text(s, desc, bx + Inches(0.18), Inches(5.3), Inches(3.6), Inches(0.9),
                 size=12, color=SLATE_600)
        bx += Inches(4.25)


def slide_data(prs):
    """Slide 4 — Dataset & Feature Engineering."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "03", "Dataset & Feature Engineering")

    # Left — dataset
    add_rect(s, Inches(0.6), Inches(1.45), Inches(5.5), Inches(5.6),
             fill=SLATE_50, line=SLATE_200, line_w=8000)
    add_text(s, "Dataset", Inches(0.8), Inches(1.65), Inches(5), Inches(0.4),
             size=14, bold=True, color=SLATE_900)
    dataset_items = [
        "178,267 videos after cleaning",
        "11 countries: US, UK, JP, BR, CA, DE, FR, IN, KR, MX, RU",
        "15 YouTube categories",
        "Source: YouTube Trending 2026 (Kaggle)",
        "category_id mapped to human labels",
        "like_ratio & comment_ratio clipped to [0, 1]",
    ]
    bullet_box(s, dataset_items, Inches(0.8), Inches(2.15), Inches(5.1), Inches(3.5),
               size=13, color=SLATE_700, accent=ACCENT)

    # Right — features
    add_rect(s, Inches(6.7), Inches(1.45), Inches(6.2), Inches(5.6),
             fill=SLATE_50, line=SLATE_200, line_w=8000)
    add_text(s, "Engineered Features (11 total)", Inches(6.9), Inches(1.65),
             Inches(5.8), Inches(0.4), size=14, bold=True, color=SLATE_900)

    groups = [
        ("NLP — Title",   ["title_length", "title_word_count", "title_caps_ratio",
                           "title_has_number", "title_has_question", "title_has_exclamation"]),
        ("Metadata",      ["number_of_tags", "publish_hour", "publish_dow"]),
        ("Categorical",   ["category", "country"]),
    ]
    gy = Inches(2.15)
    for grp, feats in groups:
        add_text(s, grp, Inches(6.9), gy, Inches(2.5), Inches(0.3),
                 size=11, bold=True, color=ACCENT)
        gy += Inches(0.3)
        for f in feats:
            add_text(s, f"  {f}", Inches(6.9), gy, Inches(5.5), Inches(0.28),
                     size=12, color=SLATE_700)
            gy += Inches(0.27)
        gy += Inches(0.1)


def slide_model(prs):
    """Slide 5 — Model architecture."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "04", "Model Architecture")

    add_text(s, "HistGradientBoostingClassifier  +  CalibratedClassifierCV",
             Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             size=22, bold=True, color=SLATE_900)
    add_text(s,
             "Why Histogram Gradient Boosting over Random Forest?  "
             "Native handling of categoricals, faster training, higher AUC on tabular data.",
             Inches(0.6), Inches(2.05), Inches(12), Inches(0.5),
             size=14, color=SLATE_600)

    # Two columns
    # Left — hyperparams
    add_rect(s, Inches(0.6), Inches(2.75), Inches(5.8), Inches(3.8),
             fill=SLATE_50, line=SLATE_200, line_w=8000)
    add_text(s, "Hyperparameters", Inches(0.8), Inches(2.95), Inches(5.3), Inches(0.4),
             size=14, bold=True, color=SLATE_900)
    params = [
        "max_iter = 300  (early stopping on 10% val split)",
        "learning_rate = 0.08",
        "max_depth = 8",
        "min_samples_leaf = 30",
        "l2_regularization = 1.0",
        "random_state = 42",
    ]
    bullet_box(s, params, Inches(0.8), Inches(3.45), Inches(5.3), Inches(2.8),
               size=13, color=SLATE_700, accent=ACCENT, dot="-")

    # Right — calibration
    add_rect(s, Inches(7.0), Inches(2.75), Inches(5.8), Inches(3.8),
             fill=SLATE_50, line=SLATE_200, line_w=8000)
    add_text(s, "Calibration (Isotonic Regression)", Inches(7.2), Inches(2.95),
             Inches(5.3), Inches(0.4), size=14, bold=True, color=SLATE_900)
    add_text(s,
             "Raw classifier scores are monotonically transformed using "
             "isotonic regression inside a 3-fold CalibratedClassifierCV. "
             "This ensures the output percentage is a statistically valid "
             "probability, not just a ranking score.\n\n"
             "Target: views >= 75th percentile  =  'viral'  (25% positive class)",
             Inches(7.2), Inches(3.45), Inches(5.3), Inches(2.8),
             size=13, color=SLATE_600)


def slide_results(prs):
    """Slide 6 — Model results."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "05", "Model Results")

    # 4 KPIs
    kpi_card(s, "Test ROC AUC",  "0.853", "held-out 20% (35,654 videos)",
             Inches(0.6), Inches(1.5), val_color=ACCENT)
    kpi_card(s, "CV ROC AUC",   "0.797", "5-fold  ±  0.005",
             Inches(3.75), Inches(1.5))
    kpi_card(s, "PR AUC",       "0.690", "average precision score",
             Inches(6.9), Inches(1.5))
    kpi_card(s, "Brier Score",  "0.124", "calibration error (lower = better)",
             Inches(10.1), Inches(1.5))

    # Feature importance bar (manual)
    features = [
        ("category",         0.1344),
        ("country",          0.1310),
        ("number_of_tags",   0.0879),
        ("title_length",     0.0563),
        ("title_word_count", 0.0551),
        ("publish_hour",     0.0534),
        ("title_caps_ratio", 0.0268),
        ("publish_dow",      0.0100),
    ]
    add_text(s, "Permutation Feature Importance (drop in ROC AUC when shuffled)",
             Inches(0.6), Inches(3.3), Inches(12), Inches(0.4),
             size=14, bold=True, color=SLATE_900)

    max_val = features[0][1]
    bar_max_w = Inches(5.5)
    fy = Inches(3.85)
    for feat, val in features:
        bar_w = bar_max_w * (val / max_val)
        add_text(s, feat, Inches(0.6), fy, Inches(2.3), Inches(0.28),
                 size=11, color=SLATE_700)
        add_rect(s, Inches(3.0), fy + Inches(0.04), bar_w, Inches(0.22),
                 fill=ACCENT if val >= 0.10 else SLATE_400)
        add_text(s, f"{val:.3f}", Inches(3.0) + bar_w + Inches(0.08),
                 fy, Inches(1), Inches(0.28), size=11, color=SLATE_600)
        fy += Inches(0.38)


def slide_app_advisor(prs):
    """Slide 7 — Viral Advisor UI."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "06", "Viral Advisor — Inputs & Score")

    # Left — inputs
    add_rect(s, Inches(0.6), Inches(1.45), Inches(5.6), Inches(5.6),
             fill=SLATE_50, line=SLATE_200, line_w=8000)
    add_text(s, "Pre-upload Inputs", Inches(0.8), Inches(1.65),
             Inches(5.1), Inches(0.4), size=14, bold=True, color=SLATE_900)
    fields = [
        ("Title",        "Free text — NLP features extracted automatically"),
        ("Category",     "15 categories (Entertainment, Music, Gaming ...)"),
        ("Country",      "11 countries in dataset"),
        ("Number of tags", "Integer 0 – 80"),
        ("Publish hour", "Slider 0 – 23 UTC"),
        ("Publish day",  "Mon – Sun dropdown"),
    ]
    fy = Inches(2.2)
    for field, desc in fields:
        add_text(s, field, Inches(0.8), fy, Inches(2.0), Inches(0.28),
                 size=12, bold=True, color=SLATE_900)
        add_text(s, desc, Inches(2.9), fy, Inches(3.1), Inches(0.28),
                 size=12, color=SLATE_600)
        fy += Inches(0.45)

    # Right — score display
    add_rect(s, Inches(6.7), Inches(1.45), Inches(6.2), Inches(2.8),
             fill=SLATE_900, line=SLATE_200, line_w=8000)
    add_text(s, "VIRALITY SCORE", Inches(6.95), Inches(1.72),
             Inches(5.5), Inches(0.35), size=9, bold=True, color=SLATE_400)
    add_text(s, "62.4%", Inches(6.9), Inches(2.1),
             Inches(5.5), Inches(0.9), size=52, bold=True, color=WHITE)
    add_text(s, "Calibrated probability of reaching top-25% views",
             Inches(6.95), Inches(3.1), Inches(5.5), Inches(0.4),
             size=12, color=SLATE_400)

    # Score bar
    add_rect(s, Inches(6.95), Inches(3.65), Inches(5.6), Inches(0.2), fill=SLATE_700)
    add_rect(s, Inches(6.95), Inches(3.65), Inches(5.6 * 0.624), Inches(0.2), fill=ACCENT)

    # Status cards
    statuses = [
        ("Assessment",  "High viral potential", SUCCESS),
        ("Confidence",  "High",                 SLATE_900),
        ("Outperforms", "62% of videos",        SLATE_900),
    ]
    sx = Inches(6.7)
    for lbl, val, col in statuses:
        add_rect(s, sx, Inches(4.35), Inches(2.0), Inches(0.85),
                 fill=WHITE, line=SLATE_200, line_w=6000)
        add_text(s, lbl, sx + Inches(0.14), Inches(4.48),
                 Inches(1.7), Inches(0.25), size=9, color=SLATE_500)
        add_text(s, val, sx + Inches(0.14), Inches(4.72),
                 Inches(1.7), Inches(0.35), size=13, bold=True, color=col)
        sx += Inches(2.07)


def slide_benchmark(prs):
    """Slide 8 — Benchmark + Recommendations."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "07", "Benchmark Comparison & Recommendations")

    # Benchmark table
    add_text(s, "Benchmark Comparison", Inches(0.6), Inches(1.45),
             Inches(7), Inches(0.4), size=16, bold=True, color=SLATE_900)
    headers = ["Metric", "Your Value", "Viral Benchmark", "Status"]
    col_w = [Inches(2.2), Inches(1.6), Inches(2.8), Inches(1.4)]
    hx = Inches(0.6)
    # Header row
    add_rect(s, Inches(0.6), Inches(2.0), sum(col_w), Inches(0.38), fill=SLATE_100)
    cx = Inches(0.6)
    for h, cw in zip(headers, col_w):
        add_text(s, h, cx + Inches(0.1), Inches(2.05), cw, Inches(0.3),
                 size=10, bold=True, color=SLATE_500)
        cx += cw
    # Rows
    rows = [
        ("Number of tags",    "15",         "10 – 22",            "in range",  SUCCESS),
        ("Title length",      "34 chars",   "34 – 66 chars",      "in range",  SUCCESS),
        ("Publish hour",      "15:00 UTC",  "13:00, 15:00, 21:00","peak hour", SUCCESS),
        ("Publish day",       "Friday",     "Thu, Fri",           "peak day",  SUCCESS),
        ("Viral keywords",    "2 hits",     ">= 1 recommended",   "hit",       SUCCESS),
    ]
    ry = Inches(2.38)
    for metric, val, bench, status, col in rows:
        add_rect(s, Inches(0.6), ry, sum(col_w), Inches(0.42),
                 fill=WHITE, line=SLATE_100, line_w=4000)
        vals_row = [metric, val, bench, status]
        cx = Inches(0.6)
        for i, (v, cw) in enumerate(zip(vals_row, col_w)):
            c = col if i == 3 else (SLATE_900 if i == 0 else SLATE_700)
            bold = i in (0, 3)
            add_text(s, v, cx + Inches(0.1), ry + Inches(0.08),
                     cw - Inches(0.1), Inches(0.32), size=12, color=c, bold=bold)
            cx += cw
        ry += Inches(0.42)

    # Recommendations
    add_text(s, "Recommendations", Inches(8.7), Inches(1.45),
             Inches(4.5), Inches(0.4), size=16, bold=True, color=SLATE_900)
    recs = [
        "Title length (34 chars) is already within the viral range of 34-66 characters.",
        "Tag count of 15 matches the viral range (10-22). Well optimized.",
        "Friday 15:00 UTC is a peak window. Optimal scheduling.",
        "2 viral keywords detected in title. Lift scores x2.1 and x1.8.",
        "Aim for like-ratio >= 2.3% and comment-ratio >= 0.085% post-upload.",
    ]
    ry = Inches(2.0)
    for r in recs:
        add_rect(s, Inches(8.7), ry, Inches(4.5), Inches(0.75),
                 fill=WHITE, line=SLATE_200, line_w=6000)
        add_rect(s, Inches(8.7), ry, Inches(0.06), Inches(0.75), fill=ACCENT)
        add_text(s, r, Inches(8.85), ry + Inches(0.1), Inches(4.2), Inches(0.62),
                 size=11, color=SLATE_700)
        ry += Inches(0.87)


def slide_nlp(prs):
    """Slide 9 — NLP keyword analysis."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "08", "NLP — Viral Keyword Analysis")

    add_text(s,
             "Lift score = P(word | viral) / P(word)",
             Inches(0.6), Inches(1.45), Inches(12), Inches(0.45),
             size=20, bold=True, color=SLATE_900)
    add_text(s,
             "For every word appearing >= 50 times, we compute how much more likely "
             "it is to appear in the top-25% videos versus the full dataset (Laplace-smoothed).",
             Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
             size=14, color=SLATE_600)

    # Lift bars
    top_kws = [
        ("trailer",    2.8),
        ("official",   2.5),
        ("new",        2.3),
        ("best",       2.1),
        ("highlights", 2.0),
        ("2026",       1.9),
        ("full",       1.85),
        ("vs",         1.8),
    ]
    add_text(s, "Top viral keywords (by lift)", Inches(0.6), Inches(2.7),
             Inches(6), Inches(0.4), size=14, bold=True, color=SLATE_900)
    bar_max = Inches(4.0)
    ky = Inches(3.2)
    for word, lift in top_kws:
        bw = bar_max * (lift / 3.0)
        add_text(s, word, Inches(0.6), ky, Inches(1.6), Inches(0.28),
                 size=12, color=SLATE_700)
        add_rect(s, Inches(2.3), ky + Inches(0.04), bw, Inches(0.22), fill=ACCENT)
        add_text(s, f"x{lift:.1f}", Inches(2.3) + bw + Inches(0.1),
                 ky, Inches(0.6), Inches(0.28), size=11, color=SLATE_600)
        ky += Inches(0.35)

    # Chip mockup
    add_rect(s, Inches(7.5), Inches(2.7), Inches(5.5), Inches(4.4),
             fill=SLATE_50, line=SLATE_200, line_w=8000)
    add_text(s, "In-app keyword chips", Inches(7.7), Inches(2.9),
             Inches(5), Inches(0.4), size=13, bold=True, color=SLATE_900)
    add_text(s, "Detected in title (highlighted blue):",
             Inches(7.7), Inches(3.35), Inches(5), Inches(0.35),
             size=11, color=SLATE_500)
    chips_hit = ["best  x2.1", "moments  x1.7"]
    cx_c = Inches(7.7)
    for chip in chips_hit:
        add_rect(s, cx_c, Inches(3.8), Inches(1.6), Inches(0.38),
                 fill=RGBColor(0xdb, 0xea, 0xfe), line=RGBColor(0xbf, 0xdb, 0xfe), line_w=6000)
        add_text(s, chip, cx_c + Inches(0.1), Inches(3.87),
                 Inches(1.4), Inches(0.28), size=11, color=RGBColor(0x1e, 0x40, 0xaf), bold=True)
        cx_c += Inches(1.75)
    add_text(s, "Suggestions from viral pool:",
             Inches(7.7), Inches(4.35), Inches(5), Inches(0.35),
             size=11, color=SLATE_500)
    chips_sug = ["trailer  x2.8", "official  x2.5", "highlights  x2.0", "2026  x1.9", "full  x1.85"]
    cx_c = Inches(7.7)
    row_c = 0
    for i, chip in enumerate(chips_sug):
        if cx_c + Inches(1.8) > Inches(12.8):
            cx_c = Inches(7.7)
            row_c += 1
        add_rect(s, cx_c, Inches(4.8) + row_c * Inches(0.5), Inches(1.7), Inches(0.38),
                 fill=WHITE, line=SLATE_200, line_w=6000)
        add_text(s, chip, cx_c + Inches(0.1), Inches(4.87) + row_c * Inches(0.5),
                 Inches(1.5), Inches(0.28), size=11, color=SLATE_700)
        cx_c += Inches(1.82)


def slide_dashboard(prs):
    """Slide 10 — Dashboard."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "09", "Dashboard — Data Insights")

    charts = [
        ("Top Categories",        "Pets & Animals leads\nwith 1.2M median views",  ACCENT),
        ("Views by Publish Hour", "13:00, 15:00, 21:00 UTC\nare peak windows",     ACCENT),
        ("Day of Week",           "Thu & Fri consistently\noutperform weekends",    ACCENT),
        ("Tag Distribution",      "Viral videos cluster\nat 10-22 tags",            ACCENT),
        ("Title Length",          "34-66 chars is\nthe sweet spot",                 ACCENT),
        ("Like Ratio",            "Music & Gaming have\nhighest engagement",        ACCENT),
    ]
    cx, cy = Inches(0.45), Inches(1.5)
    for i, (title, insight, col) in enumerate(charts):
        if i == 3:
            cx = Inches(0.45)
            cy = Inches(4.3)
        add_rect(s, cx, cy, Inches(4.1), Inches(2.55),
                 fill=SLATE_50, line=SLATE_200, line_w=6000)
        add_rect(s, cx, cy, Inches(4.1), Inches(0.06), fill=col)
        add_text(s, title, cx + Inches(0.15), cy + Inches(0.18),
                 Inches(3.8), Inches(0.38), size=13, bold=True, color=SLATE_900)
        # Mini bar placeholder
        add_rect(s, cx + Inches(0.15), cy + Inches(0.7), Inches(3.5), Inches(1.1),
                 fill=SLATE_100, line=SLATE_200, line_w=4000)
        add_text(s, "[ Chart ]", cx + Inches(0.15), cy + Inches(0.88),
                 Inches(3.5), Inches(0.5), size=10, color=SLATE_400,
                 align=PP_ALIGN.CENTER)
        add_text(s, insight, cx + Inches(0.15), cy + Inches(1.95),
                 Inches(3.8), Inches(0.52), size=11, color=SLATE_600, italic=True)
        cx += Inches(4.42)


def slide_summary(prs):
    """Slide 11 — Key results summary."""
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)
    section_header(s, "10", "Key Results")

    # Big stat strip
    stats = [
        ("178,267", "Videos Analyzed"),
        ("0.853",   "Test ROC AUC"),
        ("11",      "Countries"),
        ("15",      "Categories"),
    ]
    sx = Inches(0.6)
    for val, lbl in stats:
        add_rect(s, sx, Inches(1.6), Inches(2.9), Inches(1.6),
                 fill=SLATE_900, line=SLATE_200, line_w=6000)
        add_text(s, val, sx + Inches(0.15), Inches(1.75),
                 Inches(2.6), Inches(0.9), size=36, bold=True, color=WHITE)
        add_text(s, lbl, sx + Inches(0.15), Inches(2.65),
                 Inches(2.6), Inches(0.4), size=11, color=SLATE_400)
        sx += Inches(3.18)

    # Achievements
    add_text(s, "What we built", Inches(0.6), Inches(3.55),
             Inches(12), Inches(0.45), size=18, bold=True, color=SLATE_900)
    achievements = [
        "End-to-end data pipeline: 11 raw CSVs -> cleaned 178K-row dataset with 11 engineered features",
        "Calibrated HistGradientBoosting model achieving ROC AUC 0.853 (5-fold CV: 0.797 +/- 0.005)",
        "NLP keyword lift analysis over 178K titles — viral word signals with Laplace smoothing",
        "Production-quality Streamlit app: Viral Advisor, Dashboard, and Model Insights pages",
        "Transparent predictions: benchmark table, calibration curve, permutation feature importance",
    ]
    bullet_box(s, achievements, Inches(0.6), Inches(4.1), Inches(12.5), Inches(3.0),
               size=14, color=SLATE_700, accent=ACCENT)


def slide_closing(prs):
    """Slide 12 — Closing."""
    s = blank_slide(prs)
    fill_bg(s, SLATE_900)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=ACCENT)

    add_text(s, "Thank You",
             Inches(1.0), Inches(1.8), Inches(11), Inches(1.4),
             size=56, bold=True, color=WHITE)
    add_rect(s, Inches(1.0), Inches(3.4), Inches(2.5), Inches(0.06), fill=ACCENT)
    add_text(s,
             "YouTube Viral Advisor  —  Turning pre-upload guesswork into data.",
             Inches(1.0), Inches(3.65), Inches(11), Inches(0.55),
             size=18, color=SLATE_400)

    add_text(s, "Code & Submission",
             Inches(1.0), Inches(4.6), Inches(5), Inches(0.4),
             size=13, bold=True, color=SLATE_400)
    add_text(s,
             "analysis.py  |  model.py  |  app.py\n"
             "generate_slides.py  |  models/model_metrics.json",
             Inches(1.0), Inches(5.1), Inches(10), Inches(0.8),
             size=13, color=SLATE_500)


# ── Build ─────────────────────────────────────────────────────────────────────

def build(out: str = "rq2_project/SLIDES.pptx") -> None:
    prs = new_prs()
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_data(prs)
    slide_model(prs)
    slide_results(prs)
    slide_app_advisor(prs)
    slide_benchmark(prs)
    slide_nlp(prs)
    slide_dashboard(prs)
    slide_summary(prs)
    slide_closing(prs)
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides -> {out}")


if __name__ == "__main__":
    build()
